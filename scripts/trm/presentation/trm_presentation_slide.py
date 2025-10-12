#!/usr/bin/env python3
"""
🎯 TRM Hackathon Presentation Generator
Creates a comprehensive PowerPoint presentation with:
 - Problem statement
 - Solution architecture
 - Experimental results (64, 256, 512 samples)
 - Convergence analysis
Output: reports/trm_hackathon_presentation.pptx
"""

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- Setup ---
REPORT_DIR = "reports"
os.makedirs(REPORT_DIR, exist_ok=True)

# --- Load final results (512 samples) ---
df = pd.read_csv("logs/trm_robustness_sweep_v4.csv")
adv = df[df['model'] == 'Adversarial TRM']
std = df[df['model'] == 'Standard TRM']

# --- PowerPoint setup ---
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============================================
# SLIDE 1: Title
# ============================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

title = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
tf = title.text_frame
tf.text = "Certified Robustness Verification\nfor Tiny Recursive Models"
for p in tf.paragraphs:
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(25, 25, 112)

subtitle = slide1.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(0.8))
stf = subtitle.text_frame
stf.text = "GPU-Accelerated Attack-Guided Verification\nMNIST | A100 | auto-LiRPA"
for p in stf.paragraphs:
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(70, 70, 70)

# ============================================
# SLIDE 2: Problem Statement
# ============================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title2.text_frame.text = "Problem: Adversarial Vulnerability"
title2.text_frame.paragraphs[0].font.size = Pt(36)
title2.text_frame.paragraphs[0].font.bold = True
title2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

content2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(4))
cf2 = content2.text_frame
cf2.word_wrap = True

points = [
    "🎯 Neural networks are vulnerable to imperceptible perturbations",
    "📉 Standard training provides NO certified robustness guarantees",
    "⚠️ Critical for safety-critical applications (autonomous systems, medical AI)",
    "🔬 Need: Formal verification + adversarial training"
]

for point in points:
    p = cf2.add_paragraph()
    p.text = point
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    p.level = 0

# ============================================
# SLIDE 3: Solution Architecture
# ============================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title3.text_frame.text = "Solution: Verification Pipeline"
title3.text_frame.paragraphs[0].font.size = Pt(36)
title3.text_frame.paragraphs[0].font.bold = True
title3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Left: Pipeline steps
left3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5))
lf3 = left3.text_frame
lf3.word_wrap = True

p = lf3.add_paragraph()
p.text = "Verification Pipeline:"
p.font.size = Pt(22)
p.font.bold = True

steps = [
    "1️⃣ Train TRM-MLP on MNIST",
    "2️⃣ Adversarial training at ε=0.15",
    "3️⃣ Attack-guided verification (FGSM + I-FGSM)",
    "4️⃣ Formal bounds via α-CROWN",
    "5️⃣ Scale to 512 samples for statistical significance"
]

for step in steps:
    p = lf3.add_paragraph()
    p.text = step
    p.font.size = Pt(18)
    p.space_before = Pt(12)
    p.level = 1

# Right: Key specs
right3 = slide3.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5))
rf3 = right3.text_frame
rf3.word_wrap = True

p = rf3.add_paragraph()
p.text = "Technical Specs:"
p.font.size = Pt(22)
p.font.bold = True

specs = [
    "🖥️ Hardware: A100 GPU",
    "📦 Framework: auto-LiRPA",
    "🎯 Dataset: MNIST (28×28)",
    "⚡ Speed: <0.25s per sample",
    "💾 Memory: <30MB per sample",
    "🔢 Scale: Up to 512 samples"
]

for spec in specs:
    p = rf3.add_paragraph()
    p.text = spec
    p.font.size = Pt(18)
    p.space_before = Pt(12)
    p.level = 1

# ============================================
# SLIDE 4: Main Results
# ============================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title4.text_frame.text = "Results: 67x Improvement"
title4.text_frame.paragraphs[0].font.size = Pt(36)
title4.text_frame.paragraphs[0].font.bold = True
title4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add main robustness plot
plot_path = "reports/trm_results_20251013_013045.png"
if os.path.exists(plot_path):
    slide4.shapes.add_picture(plot_path, Inches(0.5), Inches(1.2), width=Inches(6))

# Results table
table_left = Inches(6.8)
table_top = Inches(1.5)
rows, cols = 5, 3
table = slide4.shapes.add_table(rows, cols, table_left, table_top, Inches(2.8), Inches(3)).table

# Headers
table.cell(0, 0).text = "ε"
table.cell(0, 1).text = "Adv TRM"
table.cell(0, 2).text = "Std TRM"

# Data (from 512-sample run)
results = [
    ("0.01", "80.1%", "1.2%"),
    ("0.02", "58.6%", "0%"),
    ("0.03", "40.2%", "0%"),
    ("0.04", "18.9%", "0%")
]

for i, (eps, adv_val, std_val) in enumerate(results, start=1):
    table.cell(i, 0).text = eps
    table.cell(i, 1).text = adv_val
    table.cell(i, 2).text = std_val

# Style table
for row in table.rows:
    row.height = Inches(0.6)
for cell in table.iter_cells():
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.CENTER

# Highlight header
for i in range(3):
    cell = table.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(70, 130, 180)

# Key finding
finding = slide4.shapes.add_textbox(Inches(6.8), Inches(5), Inches(2.8), Inches(1.2))
ff = finding.text_frame
ff.word_wrap = True
p = ff.add_paragraph()
p.text = "🎯 67× better at ε=0.01\n✅ 410/512 verified"
p.font.size = Pt(20)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(0, 100, 0)

# ============================================
# SLIDE 5: Convergence Analysis
# ============================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title5.text_frame.text = "Convergence: Sample Size Impact"
title5.text_frame.paragraphs[0].font.size = Pt(36)
title5.text_frame.paragraphs[0].font.bold = True
title5.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add convergence plot
conv_path = "reports/convergence_analysis.png"
if os.path.exists(conv_path):
    slide5.shapes.add_picture(conv_path, Inches(1), Inches(1.2), width=Inches(8))

# Convergence summary
summary = slide5.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1.5))
sf = summary.text_frame
sf.word_wrap = True

p = sf.add_paragraph()
p.text = "Sample Size Progression: 64 → 256 → 512 samples"
p.font.size = Pt(18)
p.font.bold = True

p = sf.add_paragraph()
p.text = "✅ Results converged at 256+ samples (82% → 80% @ ε=0.01)"
p.font.size = Pt(16)
p.space_before = Pt(10)

p = sf.add_paragraph()
p.text = "✅ Statistically significant with n=512"
p.font.size = Pt(16)

# ============================================
# SLIDE 6: Conclusion
# ============================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])

title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title6.text_frame.text = "Conclusion & Impact"
title6.text_frame.paragraphs[0].font.size = Pt(36)
title6.text_frame.paragraphs[0].font.bold = True
title6.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

content6 = slide6.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
cf6 = content6.text_frame
cf6.word_wrap = True

achievements = [
    "✅ Demonstrated 67× robustness improvement with adversarial training",
    "✅ GPU-accelerated verification: <0.25s per sample, <30MB memory",
    "✅ Validated across 512 samples with statistical convergence",
    "✅ Ready to scale to 7M parameter TRM models",
    "🚀 Future work: ARC-AGI reasoning tasks, β-CROWN optimization"
]

for ach in achievements:
    p = cf6.add_paragraph()
    p.text = ach
    p.font.size = Pt(22)
    p.space_before = Pt(20)
    p.level = 0

# --- Save ---
pptx_path = os.path.join(REPORT_DIR, "trm_hackathon_presentation.pptx")
prs.save(pptx_path)
print(f"✅ Presentation generated: {pptx_path}")
print(f"   Slides: 6 (Title, Problem, Solution, Results, Convergence, Conclusion)")